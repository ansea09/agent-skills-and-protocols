#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path
import argparse
import difflib
import json
import os
import shutil
import struct
import subprocess
import sys
import tempfile
from zipfile import ZIP_DEFLATED, ZipFile


SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
EXPECTED_DIR = SKILL_DIR / "tests" / "regression" / "expected"


def maybe_reexec_in_core_venv() -> None:
    if os.environ.get("DOC_TO_MD_REGRESSION_IN_VENV") == "1":
        return

    codex_home = Path(os.environ.get("CODEX_HOME", str(Path.home() / ".codex")))
    python = Path(
        os.environ.get(
            "MARKITDOWN_PYTHON",
            str(codex_home / "tools" / "markitdown-core-venv" / "bin" / "python"),
        )
    )
    if not python.exists() or not os.access(python, os.X_OK):
        return
    if Path(sys.executable) == python:
        return

    os.environ["DOC_TO_MD_REGRESSION_IN_VENV"] = "1"
    os.execv(str(python), [str(python), str(Path(__file__).resolve()), *sys.argv[1:]])


maybe_reexec_in_core_venv()


def find_mdown() -> str:
    env = os.environ.get("MDOWN_BIN")
    if env:
        return env
    found = shutil.which("mdown")
    if found:
        return found
    fallback = Path.home() / ".local" / "bin" / "mdown"
    if fallback.exists():
        return str(fallback)
    raise SystemExit("regression: could not find mdown")


def write_pdf(path: Path, text: str) -> None:
    objects: list[bytes] = []
    stream = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")
    objects.append(b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>")
    objects.append(
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
    )
    objects.append(b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream")
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(body))
        body.extend(f"{index} 0 obj\n".encode("ascii"))
        body.extend(obj)
        body.extend(b"\nendobj\n")
    xref_offset = len(body)
    body.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        body.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    body.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode(
            "ascii"
        )
    )
    path.write_bytes(bytes(body))


def write_docx(path: Path, text: str) -> None:
    with ZipFile(path, "w", ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>
""",
        )
        zf.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>
""",
        )
        zf.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>{text}</w:t></w:r></w:p>
  </w:body>
</w:document>
""",
        )


def write_xlsx(path: Path, text: str) -> None:
    try:
        from openpyxl import Workbook
    except ImportError as exc:
        raise SystemExit("regression: openpyxl is required in the core runtime") from exc
    wb = Workbook()
    ws = wb.active
    ws["A1"] = text
    wb.save(path)


def biff_record(opcode: int, data: bytes = b"") -> bytes:
    return struct.pack("<HH", opcode, len(data)) + data


def write_xls(path: Path, text: str) -> None:
    rows = ["Header", text]
    data = bytearray()
    data += biff_record(0x0009, struct.pack("<HH", 0x0000, 0x0010))
    data += biff_record(0x0000, b"\x00\x00" + struct.pack("<H", len(rows)) + b"\x00\x00" + struct.pack("<H", 1))
    for row_index, value in enumerate(rows):
        encoded = value.encode("latin-1")
        data += biff_record(
            0x0004,
            struct.pack("<HH", row_index, 0) + b"\x00\x00\x00" + bytes([len(encoded)]) + encoded,
        )
    data += biff_record(0x000A)
    path.write_bytes(bytes(data))


def write_pptx(path: Path, text: str) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("regression: python-pptx is required in the core runtime") from exc
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(914400, 914400, 5486400, 914400)
    textbox.text_frame.text = text
    prs.save(path)


def write_zip(path: Path, text: str) -> None:
    with ZipFile(path, "w", ZIP_DEFLATED) as zf:
        zf.writestr("inside.txt", f"{text}\n")


def write_epub(path: Path, text: str) -> None:
    with ZipFile(path, "w", ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", "application/epub+zip")
        zf.writestr(
            "META-INF/container.xml",
            """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>
""",
        )
        zf.writestr(
            "OPS/content.opf",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<package version="3.0" xmlns="http://www.idpf.org/2007/opf" unique-identifier="bookid">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:title>{text}</dc:title>
    <dc:creator>Doc To Md Regression</dc:creator>
    <dc:language>en</dc:language>
    <dc:identifier id="bookid">epub-smoke</dc:identifier>
  </metadata>
  <manifest>
    <item id="chapter1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine>
    <itemref idref="chapter1"/>
  </spine>
</package>
""",
        )
        zf.writestr(
            "OPS/chapter1.xhtml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<html xmlns="http://www.w3.org/1999/xhtml">
  <head><title>{text}</title></head>
  <body><h1>{text}</h1><p>{text} body.</p></body>
</html>
""",
        )


def fixtures() -> dict[str, tuple[str, callable]]:
    return {
        "sample.html": ("html.md", lambda path: path.write_text("<h1>HTML Smoke</h1>\n<p>Snapshot body.</p>\n", encoding="utf-8")),
        "sample.pdf": ("pdf.md", lambda path: write_pdf(path, "PDF Smoke")),
        "sample.docx": ("docx.md", lambda path: write_docx(path, "DOCX Smoke")),
        "sample.xls": ("xls.md", lambda path: write_xls(path, "XLS Smoke")),
        "sample.xlsx": ("xlsx.md", lambda path: write_xlsx(path, "XLSX Smoke")),
        "sample.pptx": ("pptx.md", lambda path: write_pptx(path, "PPTX Smoke")),
        "sample.epub": ("epub.md", lambda path: write_epub(path, "EPUB Smoke")),
        "sample.csv": ("csv.md", lambda path: path.write_text("name,value\nCSV Smoke,1\n", encoding="utf-8")),
        "sample.json": (
            "json.md",
            lambda path: path.write_text(json.dumps({"title": "JSON Smoke", "items": [1]}, indent=2), encoding="utf-8"),
        ),
        "sample.xml": (
            "xml.md",
            lambda path: path.write_text('<?xml version="1.0"?><root><title>XML Smoke</title></root>\n', encoding="utf-8"),
        ),
        "sample.zip": ("zip.md", lambda path: write_zip(path, "ZIP Smoke")),
    }


def normalize_snapshot(text: str, source: Path) -> str:
    normalized = text.replace("\r\n", "\n").replace(str(source), "<SOURCE_PATH>")
    if not normalized.endswith("\n"):
        normalized += "\n"
    return normalized


def convert_fixture(mdown: str, source: Path, output: Path) -> str:
    proc = subprocess.run([mdown, str(source), "-o", str(output)], text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if proc.returncode != 0:
        raise SystemExit(f"regression: conversion failed for {source.name}: {proc.stderr.strip()}")
    return output.read_text(encoding="utf-8", errors="replace")


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Run the doc-to-md synthetic regression corpus.")
    parser.add_argument("--update", action="store_true", help="Update expected snapshots intentionally.")
    parser.add_argument("--expected-dir", type=Path, default=EXPECTED_DIR)
    args = parser.parse_args(argv)

    mdown = find_mdown()
    tmp = Path(tempfile.mkdtemp(prefix="doc-to-md-regression."))
    try:
        args.expected_dir.mkdir(parents=True, exist_ok=True)
        failures = 0
        for source_name, (snapshot_name, writer) in fixtures().items():
            source = tmp / source_name
            output = tmp / snapshot_name
            writer(source)
            actual = normalize_snapshot(convert_fixture(mdown, source, output), source)
            expected_path = args.expected_dir / snapshot_name
            if args.update:
                expected_path.write_text(actual, encoding="utf-8")
                print(f"[ok] updated {expected_path.relative_to(SKILL_DIR)}")
                continue
            if not expected_path.is_file():
                print(f"[fail] missing expected snapshot: {expected_path}", file=sys.stderr)
                failures += 1
                continue
            expected = expected_path.read_text(encoding="utf-8")
            if actual != expected:
                print(f"[fail] snapshot drift: {snapshot_name}", file=sys.stderr)
                for line in difflib.unified_diff(
                    expected.splitlines(keepends=True),
                    actual.splitlines(keepends=True),
                    fromfile=f"expected/{snapshot_name}",
                    tofile=f"actual/{snapshot_name}",
                ):
                    sys.stderr.write(line)
                failures += 1
            else:
                print(f"[ok] snapshot {snapshot_name}")
        if failures:
            return 1
        print("[ok] doc-to-md regression corpus passed")
        return 0
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))

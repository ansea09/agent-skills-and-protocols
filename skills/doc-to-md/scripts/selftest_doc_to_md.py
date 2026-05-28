#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path
import json
import os
import shutil
import struct
import subprocess
import sys
import tempfile
from zipfile import ZIP_DEFLATED, ZipFile


def find_mdown() -> str:
    env = os.environ.get("MDOWN_BIN")
    if env:
        return env
    found = shutil.which("mdown")
    if found:
        return found
    fallback = Path.home() / ".local" / "bin" / "mdown"
    if fallback.exists():
        return str(fallback)
    raise SystemExit("selftest: could not find mdown")


def write_pdf(path: Path, text: str) -> None:
    objects: list[bytes] = []
    stream = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")
    objects.append(b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>")
    objects.append(
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
    )
    objects.append(b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream")
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(body))
        body.extend(f"{index} 0 obj\n".encode("ascii"))
        body.extend(obj)
        body.extend(b"\nendobj\n")
    xref_offset = len(body)
    body.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        body.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    body.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode(
            "ascii"
        )
    )
    path.write_bytes(bytes(body))


def write_docx(path: Path, text: str) -> None:
    with ZipFile(path, "w", ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>
""",
        )
        zf.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>
""",
        )
        zf.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>{text}</w:t></w:r></w:p>
  </w:body>
</w:document>
""",
        )


def write_xlsx(path: Path, text: str) -> None:
    try:
        from openpyxl import Workbook
    except ImportError as exc:
        raise SystemExit("selftest: openpyxl is required in the core runtime") from exc
    wb = Workbook()
    ws = wb.active
    ws["A1"] = text
    wb.save(path)


def biff_record(opcode: int, data: bytes = b"") -> bytes:
    return struct.pack("<HH", opcode, len(data)) + data


def write_xls(path: Path, text: str) -> None:
    rows = ["Header", text]
    data = bytearray()
    data += biff_record(0x0009, struct.pack("<HH", 0x0000, 0x0010))  # BIFF2 worksheet BOF.
    data += biff_record(0x0000, b"\x00\x00" + struct.pack("<H", len(rows)) + b"\x00\x00" + struct.pack("<H", 1))
    for row_index, value in enumerate(rows):
        encoded = value.encode("latin-1")
        if len(encoded) > 255:
            raise ValueError("BIFF2 test labels must fit in one byte")
        data += biff_record(
            0x0004,
            struct.pack("<HH", row_index, 0) + b"\x00\x00\x00" + bytes([len(encoded)]) + encoded,
        )
    data += biff_record(0x000A)
    path.write_bytes(bytes(data))


def write_pptx(path: Path, text: str) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("selftest: python-pptx is required in the core runtime") from exc
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(914400, 914400, 5486400, 914400)
    textbox.text_frame.text = text
    prs.save(path)


def write_zip(path: Path, text: str) -> None:
    with ZipFile(path, "w", ZIP_DEFLATED) as zf:
        zf.writestr("inside.txt", f"{text}\n")


def run_convert(mdown: str, source: Path, output: Path, expected: str) -> None:
    proc = subprocess.run([mdown, str(source), "-o", str(output)], text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if proc.returncode != 0:
        raise SystemExit(f"selftest: conversion failed for {source.name}: {proc.stderr.strip()}")
    content = output.read_text(encoding="utf-8", errors="replace")
    if expected not in content:
        raise SystemExit(f"selftest: expected {expected!r} in {output}")
    print(f"[ok] {source.suffix or source.name} -> {output.name}")


def main() -> int:
    mdown = find_mdown()
    tmp = Path(tempfile.mkdtemp(prefix="doc-to-md-selftest."))
    keep = os.environ.get("DOC_TO_MD_KEEP_SELFTEST") == "1"
    try:
        fixtures = {
            "sample.html": ("HTML Smoke", lambda path: path.write_text("<h1>HTML Smoke</h1>\n", encoding="utf-8")),
            "sample.pdf": ("PDF Smoke", lambda path: write_pdf(path, "PDF Smoke")),
            "sample.docx": ("DOCX Smoke", lambda path: write_docx(path, "DOCX Smoke")),
            "sample.xls": ("XLS Smoke", lambda path: write_xls(path, "XLS Smoke")),
            "sample.xlsx": ("XLSX Smoke", lambda path: write_xlsx(path, "XLSX Smoke")),
            "sample.pptx": ("PPTX Smoke", lambda path: write_pptx(path, "PPTX Smoke")),
            "sample.csv": ("CSV Smoke", lambda path: path.write_text("name,value\nCSV Smoke,1\n", encoding="utf-8")),
            "sample.json": (
                "JSON Smoke",
                lambda path: path.write_text(json.dumps({"title": "JSON Smoke", "items": [1]}, indent=2), encoding="utf-8"),
            ),
            "sample.xml": (
                "XML Smoke",
                lambda path: path.write_text('<?xml version="1.0"?><root><title>XML Smoke</title></root>\n', encoding="utf-8"),
            ),
            "sample.zip": ("ZIP Smoke", lambda path: write_zip(path, "ZIP Smoke")),
        }
        for filename, (expected, writer) in fixtures.items():
            source = tmp / filename
            output = tmp / f"{Path(filename).stem}.md"
            writer(source)
            run_convert(mdown, source, output, expected)

        protected = tmp / "protected.md"
        protected.write_text("old content\n", encoding="utf-8")
        proc = subprocess.run(
            [mdown, str(tmp / "missing.pdf"), "-o", str(protected)],
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        if proc.returncode == 0:
            raise SystemExit("selftest: missing input unexpectedly succeeded")
        if protected.read_text(encoding="utf-8") != "old content\n":
            raise SystemExit("selftest: protected output was modified after failed conversion")
        print("[ok] failed conversion preserved existing -o output")

        stdout_proc = subprocess.run([mdown, str(tmp / "sample.html")], text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if stdout_proc.returncode != 0 or "HTML Smoke" not in stdout_proc.stdout:
            raise SystemExit("selftest: stdout conversion failed")
        if "stdout mode is not atomic" not in stdout_proc.stderr:
            raise SystemExit("selftest: stdout warning was not emitted")
        print("[ok] stdout mode emits non-atomic warning")

        print(f"[ok] doc-to-md selftest passed: {tmp}")
        return 0
    finally:
        if not keep:
            shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    raise SystemExit(main())
